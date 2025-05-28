from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create a presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title Slide
title_content_layout = prs.slide_layouts[1]  # Title and Content
section_header_layout = prs.slide_layouts[2]  # Section Header
two_content_layout = prs.slide_layouts[3]  # Two Content
blank_layout = prs.slide_layouts[6]  # Blank

# Helper functions
def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section_slide(title):
    slide = prs.slides.add_slide(section_header_layout)
    slide.shapes.title.text = title
    return slide

def add_content_slide(title, content_list):
    slide = prs.slides.add_slide(title_content_layout)
    slide.shapes.title.text = title
    
    content = slide.placeholders[1].text_frame
    for item in content_list:
        p = content.add_paragraph()
        p.text = item
        p.level = 0
    
    return slide

def add_two_column_slide(title, left_title, left_content, right_title, right_content):
    slide = prs.slides.add_slide(two_content_layout)
    slide.shapes.title.text = title
    
    # Left content
    left = slide.placeholders[1]
    left_tf = left.text_frame
    left_tf.text = left_title
    
    for item in left_content:
        p = left_tf.add_paragraph()
        p.text = item
        p.level = 1
    
    # Right content
    right = slide.placeholders[2]
    right_tf = right.text_frame
    right_tf.text = right_title
    
    for item in right_content:
        p = right_tf.add_paragraph()
        p.text = item
        p.level = 1
    
    return slide

# Create slides
# Title Slide
add_title_slide("AI Tools Suite", "A Collection of AI-Powered Applications")

# Introduction
add_section_slide("Introduction")

add_content_slide("Project Overview", [
    "• A comprehensive suite of AI-powered tools",
    "• Leverages modern AI technologies and APIs",
    "• Designed to solve everyday problems",
    "• Demonstrates practical applications of AI",
    "• Built with web technologies and Python"
])

add_content_slide("Technologies Used", [
    "• Frontend: HTML, CSS, JavaScript",
    "• Backend: Python, Flask",
    "• AI: Google Gemini AI",
    "• APIs: OpenWeatherMap, AQICN, AssemblyAI",
    "• Web Scraping: BeautifulSoup",
    "• Media Processing: MoviePy, yt-dlp"
])

# Project Architecture
add_section_slide("Project Architecture")

add_content_slide("System Architecture", [
    "• Central index page linking to all applications",
    "• Each application operates independently",
    "• Mix of client-side and server-side applications",
    "• Environment variables for API key management",
    "• Responsive design for all applications"
])

add_two_column_slide("Frontend & Backend Technologies", 
                    "Frontend", 
                    ["• HTML5", "• CSS3", "• JavaScript", "• Responsive Design", "• Modern UI/UX"], 
                    "Backend", 
                    ["• Python 3", "• Flask Framework", "• RESTful APIs", "• Web Scraping", "• Media Processing"])

# Individual Applications
add_section_slide("Individual Applications")

# 1. AI Weather Application
add_content_slide("AI Weather Application", [
    "• Provides current weather information",
    "• Uses OpenWeatherMap API for weather data",
    "• Uses AQICN API for air quality information",
    "• Supports location search and geolocation",
    "• AI-powered responses to weather queries",
    "• Modern glassmorphism UI design"
])

add_two_column_slide("AI Weather Application", 
                    "Key Features", 
                    ["• Current weather conditions", "• Air quality information", "• Geolocation support", "• AI-powered responses", "• Responsive design"], 
                    "Technologies", 
                    ["• JavaScript", "• OpenWeatherMap API", "• AQICN API", "• Google Gemini AI", "• Geolocation API"])

# 2. FoodScanner
add_content_slide("Barcode-based AI FoodScanner", [
    "• Scans food product barcodes",
    "• Retrieves nutritional information",
    "• Provides AI-powered health assessment",
    "• Dark/light mode toggle",
    "• Camera-based barcode scanning",
    "• Flask backend for processing"
])

add_two_column_slide("Barcode-based AI FoodScanner", 
                    "Key Features", 
                    ["• Barcode scanning", "• Nutritional analysis", "• Health assessment", "• Dark/light mode", "• Mobile-friendly design"], 
                    "Technologies", 
                    ["• Python Flask", "• ZXing library", "• BeautifulSoup", "• Open Food Facts API", "• Google Gemini AI"])

# 3. AI Journal Application
add_content_slide("AI Journal Application", [
    "• Note-taking with AI summarization",
    "• Local storage for entries",
    "• Create, read, and delete functionality",
    "• Bullet point formatting",
    "• Timestamps for entries",
    "• Modern card-based UI"
])

add_two_column_slide("AI Journal Application", 
                    "Key Features", 
                    ["• Create journal entries", "• AI summarization", "• Local storage", "• Timestamps", "• Delete functionality"], 
                    "Technologies", 
                    ["• JavaScript", "• LocalStorage API", "• Google Gemini AI", "• Modern CSS", "• Responsive design"])

# 4. YouTube Video Summarizer
add_content_slide("YouTube Video Summarizer", [
    "• Downloads YouTube videos",
    "• Extracts audio from videos",
    "• Transcribes audio to text",
    "• Generates AI-powered summaries",
    "• Simple interface for URL input",
    "• Flask backend for processing"
])

add_two_column_slide("YouTube Video Summarizer", 
                    "Key Features", 
                    ["• YouTube video processing", "• Audio extraction", "• Speech-to-text", "• AI summarization", "• Text output"], 
                    "Technologies", 
                    ["• Python Flask", "• yt-dlp", "• MoviePy", "• AssemblyAI", "• Google Gemini AI"])

# 5. AI Webpage Summarizer
add_content_slide("AI Webpage Summarizer", [
    "• Scrapes content from web pages",
    "• Extracts meaningful text",
    "• Generates AI-powered summaries",
    "• Simple interface for URL input",
    "• Flask backend for processing"
])

add_two_column_slide("AI Webpage Summarizer", 
                    "Key Features", 
                    ["• Web page scraping", "• Text extraction", "• AI summarization", "• Simple interface", "• Quick processing"], 
                    "Technologies", 
                    ["• Python Flask", "• BeautifulSoup", "• Requests library", "• Google Gemini AI", "• RESTful API"])

# Technical Implementation
add_section_slide("Technical Implementation")

add_content_slide("API Integration", [
    "• OpenWeatherMap API for weather data",
    "• AQICN API for air quality information",
    "• Google Gemini AI for natural language processing",
    "• AssemblyAI for speech-to-text transcription",
    "• Open Food Facts API for product information",
    "• Environment variables for API key management"
])

add_content_slide("Security Considerations", [
    "• API keys stored in environment variables",
    "• .env file for Python applications",
    "• env-config.js for JavaScript applications",
    "• .gitignore to prevent committing sensitive data",
    "• Fallback mechanisms for missing environment variables",
    "• HTTPS for API requests where possible"
])

add_content_slide("User Interface Design", [
    "• Modern, responsive design across all applications",
    "• Glassmorphism effects in Weather application",
    "• Dark/light mode toggle in FoodScanner",
    "• Card-based layout in Journal application",
    "• Consistent styling within each application",
    "• Mobile-friendly interfaces"
])

# Implementation Details
add_section_slide("Implementation Details")

add_content_slide("Environment Variable Management", [
    "• Centralized environment variable storage",
    "• Python applications use dotenv library",
    "• JavaScript applications use window.ENV object",
    "• Fallback mechanisms for missing variables",
    "• Secure API key management"
])

add_content_slide("Google Gemini AI Integration", [
    "• Initialization of Gemini AI client",
    "• Model selection based on requirements",
    "• Prompt engineering for specific use cases",
    "• Response processing and formatting",
    "• Error handling and fallbacks"
])

add_content_slide("Web Scraping Implementation", [
    "• HTTP requests to target websites",
    "• HTML parsing with BeautifulSoup",
    "• Content extraction and filtering",
    "• Text processing and cleaning",
    "• Error handling for failed requests"
])

# Future Enhancements
add_section_slide("Future Enhancements")

add_content_slide("Potential Improvements", [
    "• Cloud synchronization for Journal application",
    "• Enhanced error handling across all applications",
    "• Caching mechanisms for API requests",
    "• User authentication and personalization",
    "• Progressive Web App (PWA) capabilities",
    "• Expanded AI capabilities with more models"
])

add_two_column_slide("Future Enhancements", 
                    "Technical Improvements", 
                    ["• Serverless architecture", "• Database integration", "• Automated testing", "• CI/CD pipeline", "• Containerization"], 
                    "Feature Enhancements", 
                    ["• Multi-language support", "• Voice commands", "• Social sharing", "• Analytics dashboard", "• Offline functionality"])

# Challenges and Solutions
add_section_slide("Challenges and Solutions")

add_content_slide("Development Challenges", [
    "• API rate limits and quotas",
    "• Cross-browser compatibility",
    "• Mobile responsiveness",
    "• Security of API keys",
    "• Integration of multiple technologies",
    "• Performance optimization"
])

add_two_column_slide("Challenges and Solutions", 
                    "Challenges", 
                    ["• API rate limits", "• Browser compatibility", "• Security concerns", "• Performance issues", "• Integration complexity"], 
                    "Solutions", 
                    ["• Caching strategies", "• Progressive enhancement", "• Environment variables", "• Code optimization", "• Modular architecture"])

# Learning Outcomes
add_section_slide("Learning Outcomes")

add_content_slide("Technical Skills Gained", [
    "• AI and machine learning integration",
    "• API development and consumption",
    "• Full-stack web development",
    "• Security best practices",
    "• User interface design",
    "• Media processing and analysis"
])

add_content_slide("Soft Skills Developed", [
    "• Problem-solving",
    "• Project management",
    "• Documentation",
    "• User experience design",
    "• Technical communication",
    "• Time management"
])

# Conclusion
add_section_slide("Conclusion")

add_content_slide("Project Summary", [
    "• Successfully created a suite of AI-powered tools",
    "• Implemented modern technologies and best practices",
    "• Demonstrated practical applications of AI",
    "• Created intuitive user interfaces",
    "• Established a foundation for future enhancements",
    "• Gained valuable technical and soft skills"
])

add_content_slide("Thank You", [
    "• Questions?",
    "• Feedback?",
    "• Suggestions?",
    "",
    "Thank you for your attention!"
])

# Save the presentation
prs.save('d:/AiTools/AI_Tools_Presentation_Updated.pptx')
print("Presentation created successfully!")
